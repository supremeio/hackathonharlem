"""
PPTX generator.

Builds an editable deck ON the Maverx master template — slides are created from
the template's named layouts (so the logo, footer, colours and patterns come
from Maverx's own master, not redrawn) and content is added as real, editable
text boxes / tables / shapes (never images).

Every content slide receives speaker notes with all required facilitation
fields. The deck follows the fixed didactic arc, with the Maverx "before /
per-topic / after" framing.
"""

from __future__ import annotations

from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from .config import BLANK_TEMPLATE, MASTER_TEMPLATE, Color, Font, Geo, LAYOUT_MAP
from .schema import Module, SpeakerNotes, Training


def _rgb(hexv: str) -> RGBColor:
    return RGBColor(int(hexv[0:2], 16), int(hexv[2:4], 16), int(hexv[4:6], 16))


class DeckBuilder:
    def __init__(self, template: Optional[Path] = None):
        path = template or (BLANK_TEMPLATE if BLANK_TEMPLATE.exists() else MASTER_TEMPLATE)
        self.prs = Presentation(str(path))
        if not BLANK_TEMPLATE.exists():
            self._strip_slides()
        self._index_layouts()

    # -- template plumbing ------------------------------------------------ #
    def _strip_slides(self) -> None:
        RELID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        lst = self.prs.slides._sldIdLst
        for sldId in list(lst):
            try:
                self.prs.part.drop_rel(sldId.get(RELID))
            except Exception:
                pass
            lst.remove(sldId)

    def _index_layouts(self) -> None:
        self._by_master = {}
        self._by_name = {}
        for mi, master in enumerate(self.prs.slide_masters):
            for layout in master.slide_layouts:
                self._by_master[(mi, layout.name)] = layout
                self._by_name.setdefault(layout.name, layout)

    def _layout(self, role: str):
        mi, name = LAYOUT_MAP[role]
        return self._by_master.get((mi, name)) or self._by_name.get(name) \
            or self.prs.slide_layouts[0]

    def add_slide(self, role: str, notes: Optional[SpeakerNotes] = None):
        slide = self.prs.slides.add_slide(self._layout(role))
        if notes is not None:
            slide.notes_slide.notes_text_frame.text = notes.render()
        return slide

    # -- text primitives -------------------------------------------------- #
    def _textbox(self, slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Pt(0)
        tf.margin_top = tf.margin_bottom = Pt(0)
        return tb, tf

    def _run(self, para, text, size, color, bold=False, italic=False, font=Font.BODY):
        r = para.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = _rgb(color)
        return r

    def title(self, slide, text, color=Color.INDIGO):
        tb, tf = self._textbox(slide, Geo.MARGIN_L, Geo.TITLE_TOP,
                               Geo.content_w(), Geo.TITLE_H, MSO_ANCHOR.MIDDLE)
        self._run(tf.paragraphs[0], text, Font.SIZE_TITLE, color,
                  bold=True, font=Font.TITLE)
        return tb

    def subtitle_under_title(self, slide, text, color=Color.CORAL):
        tb, tf = self._textbox(slide, Geo.MARGIN_L, Geo.TITLE_TOP + Geo.TITLE_H - 0.05,
                               Geo.content_w(), 0.6)
        self._run(tf.paragraphs[0], text, Font.SIZE_SUBTITLE, color,
                  bold=True, font=Font.TITLE)
        return tb

    def bullets(self, slide, main_items: List[str], sub_map=None,
                x=None, y=None, w=None, h=None, main_size=Font.SIZE_MAIN):
        """Main points (bold, no glyph) with optional sub-bullets (• regular)."""
        x = Geo.MARGIN_L if x is None else x
        y = Geo.BODY_TOP if y is None else y
        w = Geo.content_w() if w is None else w
        h = (Geo.BODY_BOTTOM - y) if h is None else h
        tb, tf = self._textbox(slide, x, y, w, h)
        first = True
        for item in main_items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(10)
            p.line_spacing = 1.05
            self._run(p, item, main_size, Color.INDIGO, bold=True)
            for sub in (sub_map or {}).get(item, []):
                sp = tf.add_paragraph()
                sp.space_after = Pt(4)
                sp.level = 1
                self._run(sp, "•  ", Font.SIZE_SUB, Color.ORANGE, bold=True)
                self._run(sp, sub, Font.SIZE_SUB, Color.BLACK)
        return tb

    def statement(self, slide, text, y, color=Color.INDIGO):
        tb, tf = self._textbox(slide, Geo.MARGIN_L, y, Geo.content_w(), 0.9)
        self._run(tf.paragraphs[0], text, Font.SIZE_SUBSUB, color,
                  bold=True, font=Font.TITLE)
        return tb

    def paragraph(self, slide, text, x, y, w, h, size=Font.SIZE_SUB,
                  color=Color.BLACK, bold=False):
        tb, tf = self._textbox(slide, x, y, w, h)
        self._run(tf.paragraphs[0], text, size, color, bold=bold)
        return tb

    # -- composite shapes ------------------------------------------------- #
    def callout(self, slide, label, body, x, y, w, h, fill=Color.PURPLE):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(x), Inches(y), Inches(w), Inches(h))
        box.fill.solid()
        box.fill.fore_color.rgb = _rgb(fill)
        box.line.fill.background()
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.25)
        tf.margin_top = tf.margin_bottom = Inches(0.15)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        self._run(p, label, Font.SIZE_SMALL, Color.WHITE, bold=True, font=Font.TITLE)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        p2.space_before = Pt(4)
        self._run(p2, body, Font.SIZE_SUB, Color.WHITE)
        return box

    def card(self, slide, head, lead_color, lines, x, y, w, h):
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = _rgb(Color.WHITE)
        bg.line.color.rgb = _rgb(Color.LAVENDER)
        bg.line.width = Pt(1)
        bg.shadow.inherit = False
        # colored dot
        d = 0.32
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.25),
                                     Inches(y + 0.22), Inches(d), Inches(d))
        dot.fill.solid()
        dot.fill.fore_color.rgb = _rgb(lead_color)
        dot.line.fill.background()
        dot.shadow.inherit = False
        # header
        htb, htf = self._textbox(slide, x + 0.7, y + 0.18, w - 0.9, 0.45)
        self._run(htf.paragraphs[0], head, Font.SIZE_CARD_HEAD, lead_color,
                  bold=True, font=Font.TITLE)
        # body
        btb, btf = self._textbox(slide, x + 0.3, y + 0.75, w - 0.6, h - 0.9)
        for i, ln in enumerate(lines):
            p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            p.space_after = Pt(5)
            p.line_spacing = 1.0
            self._run(p, ln, Font.SIZE_SMALL, Color.BLACK)
        return bg

    def number_row(self, slide, items: List[str], y=3.2):
        """A row of numbered colour circles with captions (style-guide motif)."""
        n = len(items)
        gap = 0.5
        avail = Geo.content_w()
        d = min(1.1, (avail - gap * (n - 1)) / n)
        total = d * n + gap * (n - 1)
        x0 = Geo.MARGIN_L + (avail - total) / 2
        for i, cap in enumerate(items):
            cx = x0 + i * (d + gap)
            col = Color.SEQUENCE[i % len(Color.SEQUENCE)]
            circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx), Inches(y),
                                          Inches(d), Inches(d))
            circ.fill.solid()
            circ.fill.fore_color.rgb = _rgb(col)
            circ.line.fill.background()
            circ.shadow.inherit = False
            ctf = circ.text_frame
            ctf.word_wrap = True
            self._run(ctf.paragraphs[0], str(i + 1), 24, Color.WHITE,
                      bold=True, font=Font.TITLE)
            ctf.paragraphs[0].alignment = PP_ALIGN.CENTER
            # caption under circle
            cap_tb, cap_tf = self._textbox(slide, cx - 0.35, y + d + 0.1,
                                           d + 0.7, 0.9)
            cap_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            self._run(cap_tf.paragraphs[0], cap, Font.SIZE_SMALL, Color.INDIGO,
                      bold=True)
            if i < n - 1:  # arrow
                ar = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW, Inches(cx + d + 0.05),
                    Inches(y + d / 2 - 0.06), Inches(gap - 0.1), Inches(0.12))
                ar.fill.solid()
                ar.fill.fore_color.rgb = _rgb(Color.INDIGO)
                ar.line.fill.background()
                ar.shadow.inherit = False

    def table(self, slide, headers: List[str], rows: List[List[str]],
              x=None, y=None, w=None, col_w=None):
        x = Geo.MARGIN_L if x is None else x
        y = Geo.BODY_TOP if y is None else y
        w = Geo.content_w() if w is None else w
        nrows = len(rows) + 1
        ncols = len(headers)
        h = min(5.0, 0.45 * nrows)
        gshape = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y),
                                        Inches(w), Inches(h))
        tbl = gshape.table
        if col_w:
            for ci, cw in enumerate(col_w):
                tbl.columns[ci].width = Inches(cw)
        for ci, htext in enumerate(headers):
            c = tbl.cell(0, ci)
            c.fill.solid()
            c.fill.fore_color.rgb = _rgb(Color.INDIGO)
            self._fill_cell(c, htext, Color.WHITE, bold=True, size=Font.SIZE_SMALL)
        for ri, row in enumerate(rows, start=1):
            for ci, val in enumerate(row):
                c = tbl.cell(ri, ci)
                c.fill.solid()
                c.fill.fore_color.rgb = _rgb(Color.WHITE if ri % 2 else Color.LAVENDER)
                self._fill_cell(c, str(val), Color.BLACK,
                                bold=(ci == 0), size=Font.SIZE_SMALL)
        return tbl

    def _fill_cell(self, cell, text, color, bold=False, size=12):
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        tf = cell.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = ""
        self._run(tf.paragraphs[0], text, size, color, bold=bold)

    # -- saving ----------------------------------------------------------- #
    def save(self, path) -> str:
        self.prs.save(str(path))
        return str(path)


# --------------------------------------------------------------------------- #
# Assembly: Training -> full deck
# --------------------------------------------------------------------------- #
def build_deck(tr: Training, out_path) -> str:
    d = DeckBuilder()

    _cover(d, tr)
    _about(d, tr)
    _timetable_trainer(d, tr)
    _timetable_learner(d, tr)
    _agenda(d, tr)
    _kickoff(d, tr)

    for i, m in enumerate(tr.modules, start=1):
        _module_divider(d, m, i, len(tr.modules))
        _theory(d, m)
        _example(d, m)
        _exercise(d, m)

    _wrapup_divider(d, tr)
    _takeaways(d, tr)
    _whats_next(d, tr)
    _closing(d, tr)

    return d.save(out_path)


# ---- structural slides ---------------------------------------------------- #
def _notes(aim, time, instr, dps, link, refl, deb) -> SpeakerNotes:
    return SpeakerNotes(aim=aim, time=time, instructions=instr,
                        discussion_points=dps, link_to_reality=link,
                        reflective_question=refl, debrief=deb)


def _cover(d, tr: Training):
    n = _notes(
        f"Open the session and frame what '{tr.topic}' will let people do.",
        "2 min",
        "Welcome the room warmly. Read the title, say who it's for, and give "
        "the one-line promise of the day before moving on.",
        ["Who this is for", "The promise of the session", "Energy and tone"],
        "Name something in their week this will make easier.",
        "What do you most want to walk away able to do?",
        "Set a confident, practical tone — we build something usable today.",
    )
    s = d.add_slide("cover", n)
    d.title(s, tr.title, color=Color.WHITE)
    d.subtitle_under_title(s, tr.subtitle, color=Color.ORANGE)


def _about(d, tr: Training):
    n = _notes(
        "Align everyone on goals, outcomes, audience and ground rules.",
        "4 min",
        "Walk the four cards left to right. Keep it brisk — this is the contract "
        "for the day. Invite one person to add what they want from it.",
        ["Objectives vs outcomes", "Who it's for", "How we'll work"],
        "Tie the objective back to a concrete task they own.",
        "Which of these outcomes matters most to you?",
        "We know where we're going and why — let's start.",
    )
    s = d.add_slide("content", n)
    d.title(s, "About this session")
    cw, ch = 5.55, 2.35
    gx, gy = 0.35, 0.25
    x0, y0 = Geo.MARGIN_L, 1.6
    a = tr.about
    d.card(s, "Learning objectives", Color.PURPLE, a.learning_objectives,
           x0, y0, cw, ch)
    d.card(s, "Learning outcomes", Color.ORANGE, a.learning_outcomes,
           x0 + cw + gx, y0, cw, ch)
    d.card(s, "Target group", Color.MAGENTA, [a.target_group],
           x0, y0 + ch + gy, cw, ch)
    d.card(s, "Good to know", Color.INDIGO, a.good_to_know,
           x0 + cw + gx, y0 + ch + gy, cw, ch)


def _timetable_trainer(d, tr: Training):
    n = _notes(
        "Give the trainer the run-of-show with timings.",
        "1 min (reference)",
        "This is your map for the session. Keep an eye on the clock against "
        "these segments; it's fine to flex, but protect the exercise time.",
        ["Total timing", "Where the exercises sit", "Protected debrief time"],
        "Timeboxing keeps energy up — like ad breaks pace a show.",
        "Where might we run over, and what gives?",
        "Stick to the spine; flex within segments.",
    )
    s = d.add_slide("content", n)
    d.title(s, "Timetable — trainer view")
    rows = [[r.segment, f"{r.minutes} min", r.activity] for r in tr.trainer_timetable]
    d.table(s, ["Segment", "Time", "What happens"], rows,
            col_w=[5.2, 1.4, 5.1])


def _timetable_learner(d, tr: Training):
    n = _notes(
        "Show participants the shape of the day at a glance.",
        "1 min",
        "Quickly walk the five phases so people know the journey. Emphasise the "
        "hands-on parts — that's where the value lands.",
        ["The arc of the day", "When they'll be active", "The end goal"],
        "A good session has a rhythm: learn, see, try, commit.",
        "Which phase are you most looking forward to?",
        "Now you know the journey — let's begin.",
    )
    s = d.add_slide("content", n)
    d.title(s, "Timetable — what your day looks like")
    d.number_row(s, [p.split(" — ")[0] for p in tr.learner_timetable], y=3.3)


def _agenda(d, tr: Training):
    n = _notes(
        "Preview the topics we'll cover and in what order.",
        "2 min",
        "Read the agenda. Flag how each topic builds on the last. Ask if "
        "anything they hoped for is missing so you can address it live.",
        ["Topic order", "How topics connect", "Any gaps to flag"],
        "Each topic is a building block toward the objective.",
        "Which topic is most relevant to your current work?",
        "Clear path ahead — on to the first topic.",
    )
    s = d.add_slide("content", n)
    d.title(s, "Agenda")
    main = []
    sub = {}
    for line in tr.agenda:
        if " — " in line:
            head, tail = line.split(" — ", 1)
            main.append(head)
            sub[head] = [tail]
        else:
            main.append(line)
    d.bullets(s, main, sub)


def _kickoff(d, tr: Training):
    n = _notes(
        "Set explicit learning goals so success is shared and visible.",
        "3 min",
        "State each goal as something they'll be able to do. Ask people to "
        "silently pick the goal that matters most to them — you'll revisit it "
        "in the wrap-up.",
        ["Goals are capabilities", "Personal relevance", "We'll measure against these"],
        "Goals are the finish line we'll check ourselves against later.",
        "What would make today a success for you personally?",
        "Goals set — now we build toward them.",
    )
    s = d.add_slide("content", n)
    d.title(s, "Kick-off: what you'll be able to do")
    d.bullets(s, tr.kickoff_goals)


# ---- module slides -------------------------------------------------------- #
def _module_divider(d, m: Module, i: int, total: int):
    n = _notes(
        f"Transition into module {i} of {total}: {m.name}.",
        "30 sec",
        "Mark the shift cleanly. One sentence on why this topic matters, then "
        "move into the theory.",
        [f"We're now on {m.name}", "Why it matters next"],
        "Signposting helps people follow the thread.",
        "What do you already associate with this topic?",
        "Let's dig in.",
    )
    s = d.add_slide("section", n)
    d.title(s, f"{i}. {m.name}", color=Color.WHITE)
    d.paragraph(s, "Theory  →  Example  →  Exercise", Geo.MARGIN_L, 2.6,
                Geo.content_w(), 0.6, size=Font.SIZE_SUBSUB, color=Color.ORANGE,
                bold=True)


def _theory(d, m: Module):
    s = d.add_slide("content", m.theory.notes)
    d.title(s, m.theory.title or m.name)
    d.callout(s, "DEFINITION", m.theory.definition, Geo.MARGIN_L, 1.6,
              Geo.content_w(), 1.25, fill=Color.PURPLE)
    d.bullets(s, m.theory.points, y=3.15, h=2.4)
    if m.theory.statement:
        d.statement(s, m.theory.statement, y=Geo.BODY_BOTTOM - 0.7)


def _example(d, m: Module):
    s = d.add_slide("content", m.example.notes)
    d.title(s, m.example.title or f"{m.name} in practice")
    d.callout(s, "EXAMPLE", m.example.scenario, Geo.MARGIN_L, 1.6,
              Geo.content_w(), 1.6, fill=Color.INDIGO)
    d.bullets(s, m.example.points, y=3.5, h=2.5)


def _exercise(d, m: Module):
    s = d.add_slide("content", m.exercise.notes)
    d.title(s, m.exercise.title or f"Exercise: {m.name}")
    badge = f"{m.exercise.fmt.upper()}  ·  {m.exercise.duration_min} MIN"
    d.paragraph(s, badge, Geo.MARGIN_L, 1.55, 5.0, 0.4,
                size=Font.SIZE_SMALL, color=Color.MAGENTA, bold=True)
    steps = [f"{i}. {st}" for i, st in enumerate(m.exercise.steps, start=1)]
    d.bullets(s, steps, y=2.05, h=3.2)
    if m.exercise.debrief_questions:
        d.statement(s, "Debrief", y=Geo.BODY_BOTTOM - 1.35, color=Color.INDIGO)
        qy = Geo.BODY_BOTTOM - 0.95
        tb, tf = d._textbox(s, Geo.MARGIN_L, qy, Geo.content_w(), 0.9)
        for i, q in enumerate(m.exercise.debrief_questions):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(3)
            d._run(p, "•  ", Font.SIZE_SUB, Color.ORANGE, bold=True)
            d._run(p, q, Font.SIZE_SUB, Color.BLACK)


# ---- closing slides ------------------------------------------------------- #
def _wrapup_divider(d, tr: Training):
    n = _notes(
        "Signal the shift into consolidation and reflection.",
        "30 sec",
        "Slow the pace. Tell people we're now pulling the threads together.",
        ["We're consolidating", "Reflection is coming"],
        "The wrap-up is where learning turns into intention.",
        "What's one thing already worth keeping from today?",
        "Let's lock in what matters.",
    )
    s = d.add_slide("section", n)
    d.title(s, "Wrap-up & reflection", color=Color.WHITE)


def _takeaways(d, tr: Training):
    n = _notes(
        "Consolidate the few things that must stick.",
        "4 min",
        "Read the takeaways. For each, ask 'why does this matter?' and let the "
        "room answer. Keep it to the vital few.",
        ["The vital few", "Why each matters", "Connection to the goal"],
        "If they remember nothing else, it should be these.",
        "Which takeaway will change something for you?",
        "These are your anchors going forward.",
    )
    s = d.add_slide("content", n)
    d.title(s, "Key takeaways")
    d.bullets(s, tr.wrapup.takeaways)


def _whats_next(d, tr: Training):
    n = _notes(
        "Turn learning into committed action.",
        "3 min",
        "Walk the next steps. Ask each person to pick ONE and say it out loud or "
        "write it down — commitment sticks better when spoken.",
        ["Concrete next steps", "Personal commitment", "Keeping momentum"],
        "A goal without a next step is just a wish.",
        "What's the one step you'll take this week?",
        "You leave with a plan, not just notes.",
    )
    s = d.add_slide("content", n)
    d.title(s, "What's next")
    d.bullets(s, tr.wrapup.whats_next)


def _closing(d, tr: Training):
    refl = tr.wrapup.reflection or ["What will you do differently after today?"]
    n = _notes(
        "Close with reflection and a clear send-off.",
        "3 min",
        "Pose the reflection question(s). Give a few seconds of silence, then "
        "invite one or two shares. Thank the room and end on the objective.",
        ["Personal reflection", "One commitment", "A warm close"],
        "End on intention — the first action after the door.",
        refl[0],
        "Thank you — now go and apply it.",
    )
    s = d.add_slide("closing", n)
    d.title(s, "To close it off", color=Color.WHITE)
    tb, tf = d._textbox(s, Geo.MARGIN_L, 2.4, Geo.content_w(), 3.0)
    for i, q in enumerate(refl):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        d._run(p, "•  ", Font.SIZE_MAIN, Color.ORANGE, bold=True)
        d._run(p, q, Font.SIZE_MAIN, Color.WHITE, bold=True, font=Font.TITLE)
