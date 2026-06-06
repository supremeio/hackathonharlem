"""
Maverx Training Renderer — Person C
Takes intake.json + content JSON files, produces output.pptx + pre/post-bite markdown.
"""

import argparse
import json
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

# Maps slide_type -> layout index in maverx_master.pptx
LAYOUT_MAP = {
    "kickoff_title":         0,  # TITLE — cover
    "kickoff_agenda":        2,  # CUSTOM_1 — agenda with photo
    "kickoff_objectives":    6,  # CUSTOM_2 — title + bullets
    "theory_intro":          3,
    "theory_deep":           5,
    "example_walkthrough":   4,  # CUSTOM_3 — two-column
    "example_case":          5,
    "exercise_setup":        3,  # title + bullets (steps list)
    "exercise_instructions": 5,
    "wrapup_summary":        7,  # CUSTOM_6 — closing
    "wrapup_nextlevel":      7,
    "mentimeter_recap":      5,  # CUSTOM_4_1 — 3-card grid
}
FALLBACK_LAYOUT = 3  # title + bullets if slide_type unknown


def fill_text_placeholder(placeholder, text):
    """Replace placeholder text, preserving its formatting."""
    tf = placeholder.text_frame
    tf.text = text  # first paragraph
    # Word-wrap so long lines don't overflow
    tf.word_wrap = True


def fill_body_with_bullets(placeholder, items):
    """Fill a body placeholder with a list of bullet points."""
    tf = placeholder.text_frame
    tf.word_wrap = True
    # First item goes in the existing paragraph
    tf.text = items[0] if items else ""
    # Remaining items become new paragraphs
    for item in items[1:]:
        p = tf.add_paragraph()
        p.text = item


def render_speaker_notes(slide, notes):
    """Write the 6-field facilitation guide into the slide's speaker notes."""
    if not notes:
        return
    nf = slide.notes_slide.notes_text_frame
    lines = []
    if notes.get("aim"):
        lines.append(f"AIM — {notes['aim']}")
    if notes.get("time"):
        lines.append(f"TIME — {notes['time']}")
    if notes.get("instructions"):
        lines.append(f"INSTRUCTIONS — {notes['instructions']}")
    kdp = notes.get("key_discussion_points") or []
    if kdp:
        lines.append("KEY DISCUSSION POINTS:")
        for point in kdp:
            lines.append(f"  • {point}")
    if notes.get("link_to_reality"):
        lines.append(f"LINK TO REALITY — {notes['link_to_reality']}")
    if notes.get("reflective_question"):
        lines.append(f"REFLECTIVE QUESTION — {notes['reflective_question']}")
    if notes.get("debrief"):
        lines.append(f"DEBRIEF — {notes['debrief']}")
    nf.text = "\n".join(lines)


def get_placeholders_by_type(slide, types):
    """Return placeholders matching given type names, in idx order."""
    matched = []
    for ph in slide.placeholders:
        ph_type = str(ph.placeholder_format.type)
        if any(t in ph_type for t in types):
            matched.append(ph)
    return sorted(matched, key=lambda p: p.placeholder_format.idx)


def render_slide(prs, slide_data):
    """Add one slide to the presentation based on slide_data dict."""
    slide_type = slide_data.get("slide_type", "")
    layout_idx = LAYOUT_MAP.get(slide_type, FALLBACK_LAYOUT)

    # Safety: clamp to available layouts
    if layout_idx >= len(prs.slide_layouts):
        layout_idx = FALLBACK_LAYOUT

    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)

    title_text = slide_data.get("title", "")
    body_items = slide_data.get("body", []) or []

    # Fill title (any placeholder type containing TITLE)
    title_phs = get_placeholders_by_type(slide, ["TITLE"])
    if title_phs and title_text:
        fill_text_placeholder(title_phs[0], title_text)

    # Fill body — handle 1 or many body placeholders
    body_phs = get_placeholders_by_type(slide, ["BODY", "SUBTITLE"])
    # Skip the title placeholder if it was already counted as TITLE
    body_phs = [p for p in body_phs if "TITLE" not in str(p.placeholder_format.type)]

    if body_phs and body_items:
        if len(body_phs) == 1:
            # Single body — pour all bullets in
            fill_body_with_bullets(body_phs[0], body_items)
        else:
            # Multiple body placeholders — distribute bullets across them
            # Each placeholder gets one body item; extras go in the last one
            for i, ph in enumerate(body_phs):
                if i < len(body_items) - 1:
                    fill_text_placeholder(ph, body_items[i])
                elif i == len(body_items) - 1:
                    fill_text_placeholder(ph, body_items[i])
                else:
                    # No content left — leave placeholder empty
                    fill_text_placeholder(ph, "")

    # Speaker notes
    render_speaker_notes(slide, slide_data.get("speaker_notes"))

    return slide


def render_deck(intake_path, content_path, master_path, out_path):
    """Build the full .pptx from intake + slides JSON."""
    intake = json.loads(Path(intake_path).read_text())
    content = json.loads(Path(content_path).read_text())

    slides = content.get("slides", [])
    if not slides:
        raise ValueError(f"No slides found in {content_path}")

    prs = Presentation(master_path)

    # python-pptx opens templates with their sample slides — remove them
    # so the output starts fresh.
    xml_slides = prs.slides._sldIdLst
    for sld_id in list(xml_slides):
        xml_slides.remove(sld_id)

    for slide_data in slides:
        render_slide(prs, slide_data)

    prs.save(out_path)
    print(f"[RENDERED] {out_path}  ({len(slides)} slides)")


def render_bite_markdown(content_path, out_path, kind):
    """Render pre/post-bite JSON as a clean markdown document."""
    data = json.loads(Path(content_path).read_text())

    lines = []
    title = data.get("title") or f"{kind.capitalize()} document"
    lines.append(f"# {title}\n")

    if data.get("introduction"):
        lines.append(data["introduction"])
        lines.append("")

    sections = data.get("sections") or []
    for sec in sections:
        if sec.get("heading"):
            lines.append(f"## {sec['heading']}\n")
        if sec.get("body"):
            lines.append(sec["body"])
            lines.append("")
        items = sec.get("items") or []
        for item in items:
            lines.append(f"- {item}")
        if items:
            lines.append("")

    # Fallback: dump any top-level lists/strings the schema might emit
    if not sections:
        for key, val in data.items():
            if key in ("title", "introduction", "schema_version", "model_name",
                       "cost_tracking", "generation_mode", "output_language"):
                continue
            if isinstance(val, str) and val.strip():
                lines.append(f"## {key.replace('_', ' ').capitalize()}\n")
                lines.append(val)
                lines.append("")
            elif isinstance(val, list):
                lines.append(f"## {key.replace('_', ' ').capitalize()}\n")
                for item in val:
                    lines.append(f"- {item}")
                lines.append("")

    Path(out_path).write_text("\n".join(lines))
    print(f"[RENDERED] {out_path}")


def main():
    parser = argparse.ArgumentParser(description="Maverx training renderer")
    parser.add_argument("--intake", required=True, help="Path to intake.json")
    parser.add_argument("--content", required=True, help="Path to slides_L1.json (or slides_Ln.json)")
    parser.add_argument("--prebite", help="Path to prebite_L1.json (optional)")
    parser.add_argument("--postbite", help="Path to postbite_L1.json (optional)")
    parser.add_argument("--master", default="maverx_master.pptx", help="Path to master template")
    parser.add_argument("--output", default="output.pptx", help="Output .pptx path")
    parser.add_argument("--output-dir", default=".", help="Directory for all outputs")
    args = parser.parse_args()

    out_dir = Path(args.output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    pptx_path = out_dir / Path(args.output).name
    render_deck(args.intake, args.content, args.master, str(pptx_path))

    if args.prebite:
        render_bite_markdown(args.prebite, str(out_dir / "prebite.md"), "pre-bite")
    if args.postbite:
        render_bite_markdown(args.postbite, str(out_dir / "postbite.md"), "post-bite")


if __name__ == "__main__":
    main()
